##############################################################################################
# PPT Quiz Maker:
# This is a program that takes a PPT and allows the user to indicate which
#section(s) they deem important by using "<@" and "@>" in the PPT which this 
#program will then extract and create fill-in-the-blank questions. The blanks
#can either be randomized or customized by using "<!" and "!>" in the PPT.
#The user can input their responses and is given a score at the end.

# This code is divided into 4 main sections: the first is the construction of
#the flashcards by extracting text to create a dictionary, the second is the
#interaction with the user by asking for responses and storing responses/scores,
#the third is the main section in which the various functions are called and the
#file name and omitted words are kept, lastly is the post-interaction where the
#total scores are evaluated and final results are displayed.
##############################################################################################

import pptx
import random
import re
import numpy as np

##############################################################################################
# CONSTRUCTION OF FLASHCARD
##############################################################################################

###LINES 39 AND 43 WERE TAKEN FROM A SOURCE AND SLIGHTLY MODIFIED FOR OUR PROGRAM

def isolate(text,start,end):
  ''' (string,string,string) -> list,list
  Descrip: Any text can be added and separated. It is best to have unique
  values for the start and end so the program knows exactly what to extract.
  
  Returns: Two lists. 1st list is a list of all indexes of where the start 
  string has appeared. 2nd list shows all indexes of where the end string has 
  appeared.
  '''
  #find all instances of start symbol within text
  #store all index of all start symbols to [start]
  index_start = [m.start() for m in re.finditer(start, text)]
  
  #find all instances of end symbol
  #store all index of all end symbols to [end]
  index_end = [m.start() for m in re.finditer(end, text)]
  return index_start, index_end

def checking_indicator(s_index,e_index,slide_number):
  ''' (list, list, int) -> bool
  Descrip: Compares the length of the first list and the second list to see whether
  they have an equal amount of items.
  
  Returns: True if the length of the first list and the second list are not equal.
  If it is true, a statement is printed indicating which slide the error is on.
  '''

  if not (len(s_index) == len(e_index)):
        print("You must have an equal amount of starting and ending indicators")
        print("Error is on slide {}. Please follow proper formatting for our code! <@ to begin and @> to end\n\n".format(slide_number)) 
        return True

###LINES 63-71 AND 73-91 (EXCLUDING THE COMMENTS) WERE TAKEN FROM A SOURCE 
###AND PARTIALLY MODIFIED TO FIT OUR PROGRAMS NEEDS

def upload(name_of_file):
  ''' (str) -> obj
  Descrip: The name of the PPT file is inputted to this function and is used to 
  upload the PPT to the program
  
  Returns: Object using the built-in function from the library pptx
  '''  
  prs = pptx.Presentation(name_of_file)
  return prs

def text_extracter(prs):
  ''' (obj) -> dict
  Descrip: Given the PPT object, this function goes through each shape in the 
  slides to find a textbox which contains the string "<@". After using the isolate
  and checking_indicator functions, each section is added to the dict as a list.
  
  Returns: Dictionary which contains the isolated strings in order of
  appearance with the associated keys being integers starting from 1.
  '''
  #reset all variables
  sentences = {}
  page_count = 1
  question_count = 1
  #go through each slide
  for slide in prs.slides:
    #go through each object in slide
    for shapes in slide.shapes:
      #if object is a textbox that contains "<@" do next step, else go to next object
      if shapes.has_text_frame and "<@" in shapes.text:
        #use isolate function to create lists with start and end index values
        s_index,e_index = isolate(shapes.text,"<@","@>")
        #check to see whether both lists are equal lengths
        #if they are not equal then move on to next object
        if checking_indicator(s_index,e_index,page_count) == True:
          pass
        #else go through the lists and add each extracted section to dict
        #index values of first list and second list create one question
        #need to add 2 for first because do not want to include "<@" in question
        else:
          for r in range(0, len(s_index)):
            sentences[question_count] = [shapes.text[s_index[r]+2:e_index[r]]]
            question_count += 1
    page_count += 1
  #if dictionary is empty, then need to have a message indicating instructions
  if sentences == {}:
    print("You need to insert <@ and @> between the section you want to be a question.")
  #return the completed dictionary
  return sentences


def flash_card_construction(sentences,common_words):
  ''' (dict,list) -> dict
  Descrip: Go through each value in the dictionary to create fill in the blank
  questions while ensuring that the list of common words is omitted. These fill
  in the blank questions are added to the associated key as a list.
  
  Returns: Updated dictionary that includes the question with the blanks as a
  list and another list with the words/section blanked out. Both being added to
  the associated key.
  '''  
  #Iterating through the dictionary
  # x is the key
  # y is the value within the key. Which is a List
  for x , y in sentences.items():
    #iterate through the list that is in y
      for z in y:
        #boolean to check whether the chosen words, if applicable, could be blanked out
          chosen_works = False
          #if the string in the list that was in y has <! and !> then it'll extract portion
          if "<!" in z and "!>" in z:
            #isolates the specified section from the string by creating lists
            a, b = isolate(z,"<!","!>")
            #check to ensure the user inputted correct number of start/end indicators
            if len(a) == len(b):
              replaced_string = []
              #temp variable to store the phrase between the start and end
              modified_sentence = z 
              for i in range(len(a)):
                # For first loop, z and sentence have same length
                # once we insert "____" the character index of a and b need to update
                # That is why we take the difference to update a, b to split the word on a new point 
                take_out = modified_sentence[a[i]+2 - (len(z) -len(modified_sentence)):b[i] - (len(z) - len(modified_sentence))]
              #splits the original text BASED off the replaced word
              #make a list[words before replaced word, words after replaced word]
                words = modified_sentence.split(sep="<!" + str(take_out) + "!>")
              #since words list is length 2, we are inserting the blank between, at [1]
                words.insert(1,"____")
                #Joins the elements of the list together with a space in between
                modified_sentence = " ".join(words)
                #adds section taken out to replaced string list
                replaced_string.append(take_out)
                #indicates that the chosen section worked without error, randomization
                #of blanks is not necessary
                chosen_works = True
                #variable that will contain final sentence to be added to dict
                new = modified_sentence
            #else provide a message with instructions
            else:
              print("Please follow proper formatting for our code! <! to begin and !> to end. The error is on question " + str(x) + 
                    ". The fill in the blank is randomized while this is not corrected.\n")
              #remove any indicators they put so that it can be randomized later
              z = str(z).replace("<!", "")
              z = str(z).replace("!>", "")
          #if the user did not specify anything to be blanked out or they did so
          #incorrectly, resulting in the randomization of blanking out words
          if chosen_works == False:
            #splits the string based off space
            words = z.split()
            #reset variables
            replaced_string = []
            temp = {}
            #at least 1/4th of the string should be blanked out
            #if there are less than 4 words then 1 word will be blanked
            if len(words) < 4:
              max_randoms = 1
            else:
              max_randoms = len(words)//4
            #get a random value between index 0 and max randoms
            for i in range(0,max_randoms):
              a = random.randint(0,len(words)-1)
              #ensure that the word of the ath value is not in common words list
              #or is already a blanked out word
              while (words[a] in common_words) or words[a] == "_____":
                #if it is then choose new random value
                a = random.randint(0,len(words)-1)
              #ensure that the words taken out do not have these characters attached
              words[a]= str(words[a]).replace(".","")
              words[a]= str(words[a]).replace(",","")
              words[a]= str(words[a]).replace("!","")
              words[a]= str(words[a]).replace(";","")
              words[a]= str(words[a]).replace('"',"")
              #switch around values (word will now become a blank)
              replaced_string, words[a] = words[a],"_____"
              #temp variable to store the word(s) selected
              temp[a] = replaced_string
            #dictionary of replaced words are out of order
            #this will sort by keys  
            #first reset replaced string to an empty list
            replaced_string = []
            #goes through using a sorted keys
            for i in sorted(temp.keys()) :
              #adds to replaced string
              replaced_string.append(temp[i])
            #Joins the elements of the list together with a space in between
            new = " ".join(words)
      #We store the word/phrase with the "_____" into the list in the dictionary
      y.append(new)
      #We strore the word/phrase we took out to the list, to show answer
      y.append(replaced_string)
  #returns the updated dictionary
  return sentences



##############################################################################################
# INTERACTION WITH FLASHCARD
##############################################################################################

def interaction(flashcards):
  ''' (dict) -> dict
  Descrip: Takes dictionary previously created with flash_card_construction function
  and provides interactive questions to the user.
  
  Returns: Updated dictionary that inclucdes the answers inputted by the user
  and if SKIP is inputted then "" is added as answer for all blanks including
  and after SKIP was inputted.
  '''  
  #go through dictionary and print question to user
  for i in sorted(flashcards.keys()):
    print("\nQuestion {}:\n {}".format(i, flashcards[i][1]))
    print("To skip question, type SKIP")
    #temp list to store the responses
    temp_list =[]
    #allow user to input answers equivalent to the number of blanks
    for j in range(len(flashcards[i][2])):
      response = (input("Fill blank number {}:\n".format(j+1))) 
      #if user types SKIP then they can skip the rest of the question
      #program assumes "" as the response for the rest of the blanks
      if response == "SKIP":
        for y in range(len((flashcards[i][2])) - j):
          temp_list.append(" ")
        #move on to the next question
        break
      else:
        temp_list.append(response)
    #add the list of responses to the appropriate key
    flashcards[i].append(temp_list)
  #return updated dictionary
  return flashcards

def check_response(flashcards_response):
  ''' (dict) -> dict
  Descrip: This function compares the answers and responses to create a list with
  0s and 1s, depending on whether the reponse was correct. This list is added 
  to the associated key.
  
  Returns: Updated dictionary with list of scores for each blank for each key
  '''
  #i is the key
  for i in sorted(flashcards_response.keys()):
    #reset variable
    question_score = []
    #j goes until the length of list of answers (i.e # of blanks)
    for j in range(len(flashcards_response[i][2])):
      #if the response is the same as the answer then the blank is given a 1
      #else it is given a 0
      #have response and answer be uppercase to ensure no case sensitivity
      if flashcards_response[i][2][j].upper() == flashcards_response[i][3][j].upper():
        question_score.append(1)
      else:
        question_score.append(0)
    #add the list of scores to the appropriate key
    flashcards_response[i].append(question_score)      




##############################################################################################
# MAIN SECTION  
##############################################################################################



#If needed, common words can be expanded upon to include words from a list found online in the future
#this list can be added upon by manually inserting new words
common_words = ["is","this","the","a","i","an","as","to","then","in","do","are","and",
                "of","for", "at", "with", "has", "that", "or", "be"]
#The user needs to type in the ppt file in the following line
file_name = "Chapter-18-Lecture-Notes-8th-Edition_Testing.pptx"
#file_name can be inputted by the user following a question if desired
#file_name = input("What is the name of your ppt file?\n")



#main operators that go through the various functions to get ultimate dict
if __name__ == "__main__":
  prs = upload(file_name)
  text_extracted = text_extracter(prs)
  flashcards = flash_card_construction(text_extracted,common_words)
  flashcards_response = interaction(flashcards)
  check_response(flashcards_response)



##############################################################################################
# POST-INTERACTION WITH FLASHCARD
##############################################################################################


#reset variables
overall_score = 0
total_questions = 0
#w is the key
#l is value of the key
for w, l in flashcards.items():
  #sum the values of the score list and provide score for each question
  print("\nYour Score for Question {}".format(w))
  print(np.array(l[4]).sum())
  #show user what their responses were
  x = ",".join(flashcards[w][3])
  print("Here are your answers: {}".format(x))
  #show user what the actual answers were
  y =",".join(flashcards[w][2])
  print("Here are the answers: {}".format(y))
  #add the score and # of questions while going through each key
  overall_score += np.array(l[4]).sum()
  total_questions += len((flashcards[w][4]))

#give total score to user and calculate percentage of score,
#if there was any question to begin with
if total_questions > 0:
  percentage = overall_score/total_questions  
  print("\nYour total score is {score} out of {questions}".format(score = overall_score, questions = total_questions))   
else:
  percentage = "undefined"

#use percentage to print message
#if there are no questions there will already be an error message printed
if percentage == "undefined":
  print("")
elif percentage < 0.5:
  print("Sorry you failed. You should try again!")
elif percentage == 1:
  print("Perfect! Feel free to try again!")
elif percentage >= 0.8 and percentage < 1:
  print("Good job! Feel free to try again!")
else:
  print("Nice Try. You should try again!")